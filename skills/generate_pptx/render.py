"""渲染：python-pptx 按版式填充内容，图表使用 PPT 原生图表（可二次编辑）。

参考 PPTAgent 的布局感知思想，渲染时读取占位符实际尺寸，自适应调整内容。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt, Emu

from skills.generate_pptx.layout_parser import get_layout_by_name

# ── 字体格式配置 ──────────────────────────────────────────────
FONT_NAME_CN = "微软雅黑"          # 中文字体
FONT_NAME_EN = "Microsoft YaHei"   # 英文字体（与中文保持一致）

TITLE_FONT_SIZE = Pt(28)           # 标题字号
SUBTITLE_FONT_SIZE = Pt(18)        # 副标题字号
BODY_FONT_SIZE = Pt(16)            # 正文字号
BODY_FONT_SIZE_SMALL = Pt(14)      # 正文小字号（空间不足时降级）
CHART_TITLE_FONT_SIZE = Pt(14)     # 图表标题字号

TITLE_COLOR = RGBColor(0x1F, 0x1F, 0x1F)    # 标题颜色（深灰）
SUBTITLE_COLOR = RGBColor(0x59, 0x59, 0x59)  # 副标题颜色（中灰）
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)      # 正文颜色（深灰）

# 段落间距配置
PARA_SPACE_BEFORE = Pt(4)          # 段前间距
PARA_SPACE_AFTER = Pt(2)           # 段后间距
PARA_LINE_SPACING = 1.2           # 行距倍数

# 正文占位符高度阈值（英寸）：低于此值使用小字号
BODY_HEIGHT_THRESHOLD = 3.0
# ──────────────────────────────────────────────────────────────

# 图表类型映射
CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
}


def _apply_font(run, font_name=FONT_NAME_CN, font_size=BODY_FONT_SIZE,
                bold=False, italic=False, color=BODY_COLOR):
    """统一设置 run 的字体属性，包括拉丁、东亚、复杂脚本三种字体。"""
    font = run.font
    font.name = font_name
    font.size = font_size
    font.bold = bold
    font.italic = italic
    if color is not None:
        font.color.rgb = color
    # 设置东亚字体（a:ea）和复杂脚本字体（a:cs），确保中文字符正确渲染
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        rPr = run._r.get_or_add_rPr()
        for tag in ('a:ea', 'a:cs'):
            node = rPr.find(qn(tag))
            if node is None:
                node = etree.SubElement(rPr, qn(tag))
            node.set('typeface', font_name)
    except Exception:
        pass


def _apply_paragraph_format(paragraph, space_before=PARA_SPACE_BEFORE,
                            space_after=PARA_SPACE_AFTER,
                            line_spacing=PARA_LINE_SPACING):
    """统一设置段落格式（段前/段后间距、行距）。"""
    try:
        paragraph.space_before = space_before
        paragraph.space_after = space_after
        paragraph.line_spacing = line_spacing
    except Exception:
        # 回退：通过 XML 设置
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = paragraph._p.get_or_add_pPr()
            if space_before is not None:
                spcBef = pPr.find(qn('a:spcBef'))
                if spcBef is None:
                    spcBef = etree.SubElement(pPr, qn('a:spcBef'))
                spcPts = spcBef.find(qn('a:spcPts'))
                if spcPts is None:
                    spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
                spcPts.set('val', str(int(space_before / 12700)))
            if space_after is not None:
                spcAft = pPr.find(qn('a:spcAft'))
                if spcAft is None:
                    spcAft = etree.SubElement(pPr, qn('a:spcAft'))
                spcPts = spcAft.find(qn('a:spcPts'))
                if spcPts is None:
                    spcPts = etree.SubElement(spcAft, qn('a:spcPts'))
                spcPts.set('val', str(int(space_after / 12700)))
        except Exception:
            pass


def _find_placeholder(slide, ph_type_id: int):
    """在幻灯片中查找指定类型的占位符。"""
    for ph in slide.placeholders:
        if ph.placeholder_format.type == ph_type_id:
            return ph
    return None


def _find_placeholder_by_idx(slide, idx: int):
    """在幻灯片中查找指定索引的占位符。"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _get_placeholder_height_inches(ph) -> float:
    """获取占位符高度（英寸）。"""
    try:
        return Emu(ph.height).inches
    except Exception:
        return 0.0


def _fill_title(slide, title_text: str) -> None:
    """填充标题占位符。"""
    if not title_text:
        return
    ph = _find_placeholder(slide, PP_PLACEHOLDER.TITLE)
    if ph is None:
        ph = _find_placeholder(slide, PP_PLACEHOLDER.CENTER_TITLE)
    if ph is None:
        ph = _find_placeholder_by_idx(slide, 0)
    if ph is None:
        return

    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    _apply_font(run, font_size=TITLE_FONT_SIZE, bold=True, color=TITLE_COLOR)
    _apply_paragraph_format(p, space_before=Pt(0), space_after=Pt(0))


def _fill_subtitle(slide, subtitle_text: str) -> None:
    """填充副标题占位符。"""
    if not subtitle_text:
        return
    ph = _find_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
    if ph is None:
        return

    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle_text
    _apply_font(run, font_size=SUBTITLE_FONT_SIZE, color=SUBTITLE_COLOR)
    _apply_paragraph_format(p, space_before=Pt(0), space_after=Pt(0))


def _fill_body(slide, bullets: list[str]) -> None:
    """填充正文占位符（要点列表），根据占位符高度自适应字号。"""
    if not bullets:
        return
    ph = _find_placeholder(slide, PP_PLACEHOLDER.BODY)
    if ph is None:
        ph = _find_placeholder(slide, PP_PLACEHOLDER.OBJECT)
    if ph is None:
        ph = _find_placeholder_by_idx(slide, 1)
    if ph is None:
        return

    # 根据占位符高度决定字号
    ph_height = _get_placeholder_height_inches(ph)
    font_size = BODY_FONT_SIZE_SMALL if ph_height > 0 and ph_height < BODY_HEIGHT_THRESHOLD else BODY_FONT_SIZE

    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    # 开启自动缩小：文本过多时自动缩小字号以适配占位符
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(bullet)
        p.level = 0
        _apply_font(run, font_size=font_size, color=BODY_COLOR)
        _apply_paragraph_format(p)


def _find_chart_area(slide) -> tuple:
    """根据幻灯片占位符布局，推算图表应放置的位置和大小。

    策略：查找 body/object 占位符的位置作为图表区域；
    若无可用占位符，使用默认位置。
    """
    ph = _find_placeholder(slide, PP_PLACEHOLDER.BODY)
    if ph is None:
        ph = _find_placeholder(slide, PP_PLACEHOLDER.OBJECT)
    if ph is not None:
        try:
            return ph.left, ph.top, ph.width, ph.height
        except Exception:
            pass
    # 默认位置
    return Inches(1.0), Inches(2.0), Inches(8.0), Inches(4.5)


def _add_chart(slide, chart_data: dict, left=None, top=None, width=None, height=None) -> None:
    """在幻灯片中添加原生图表，位置自适应占位符。"""
    chart_type_str = chart_data.get("type", "bar")
    chart_type = CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    categories = chart_data.get("categories", [])
    series_list = chart_data.get("series", [])

    data = CategoryChartData()
    data.categories = categories
    for s in series_list:
        data.add_series(s.get("name", ""), s.get("values", []))

    # 如果未指定位置，尝试从占位符推算
    if left is None or top is None or width is None or height is None:
        area_left, area_top, area_width, area_height = _find_chart_area(slide)
        if left is None:
            left = area_left
        if top is None:
            top = area_top
        if width is None:
            width = area_width
        if height is None:
            height = area_height

    chart_frame = slide.shapes.add_chart(
        chart_type, left, top, width, height, data
    )
    chart = chart_frame.chart
    chart.has_legend = True

    # 设置图表标题
    chart_title = chart_data.get("title", "")
    if chart_title:
        chart.has_title = True
        title_frame = chart.chart_title.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        run = p.add_run()
        run.text = chart_title
        _apply_font(run, font_size=CHART_TITLE_FONT_SIZE, bold=True, color=TITLE_COLOR)

    # 设置图例字体
    try:
        if chart.has_legend:
            for p in chart.legend.text_frame.paragraphs:
                for run in p.runs:
                    _apply_font(run, font_size=Pt(10), color=BODY_COLOR)
    except Exception:
        pass

    # 设置坐标轴刻度标签字体（通过 XML 操作）
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        for axis in [chart.value_axis, chart.category_axis]:
            if axis is None:
                continue
            # 设置轴标题字体
            if axis.has_title:
                for p in axis.axis_title.text_frame.paragraphs:
                    for run in p.runs:
                        _apply_font(run, font_size=Pt(10), color=BODY_COLOR)
            # 设置刻度标签字体
            try:
                txPr = axis._element.find(qn('c:txPr'))
                if txPr is None:
                    txPr = etree.SubElement(axis._element, qn('c:txPr'))
                bodyPr = txPr.find(qn('a:bodyPr'))
                if bodyPr is None:
                    bodyPr = etree.SubElement(txPr, qn('a:bodyPr'))
                lstStyle = txPr.find(qn('a:lstStyle'))
                if lstStyle is None:
                    lstStyle = etree.SubElement(txPr, qn('a:lstStyle'))
                # 添加默认段落和 run 设置字体
                p_elem = txPr.find(qn('a:p'))
                if p_elem is None:
                    p_elem = etree.SubElement(txPr, qn('a:p'))
                r_elem = p_elem.find(qn('a:r'))
                if r_elem is None:
                    r_elem = etree.SubElement(p_elem, qn('a:r'))
                rPr = r_elem.find(qn('a:rPr'))
                if rPr is None:
                    rPr = etree.SubElement(r_elem, qn('a:rPr'))
                rPr.set('lang', 'zh-CN')
                rPr.set('dirty', '0')
                rPr.set('sz', '800')  # 8pt = 800 hundredths of pt
                # 设置拉丁和东亚字体
                latin = rPr.find(qn('a:latin'))
                if latin is None:
                    latin = etree.SubElement(rPr, qn('a:latin'))
                latin.set('typeface', FONT_NAME_CN)
                ea = rPr.find(qn('a:ea'))
                if ea is None:
                    ea = etree.SubElement(rPr, qn('a:ea'))
                ea.set('typeface', FONT_NAME_CN)
            except Exception:
                pass
    except Exception:
        pass


def _fill_notes(slide, notes_text: str) -> None:
    """填充备注栏讲稿。"""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def render_slide(prs: Presentation, page: dict, layout_name: str | None = None) -> None:
    """渲染单页幻灯片。"""
    target_layout_name = layout_name or page.get("layout_name", "")

    layout = get_layout_by_name(prs, target_layout_name)
    if layout is None:
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)

    # 填充标题
    _fill_title(slide, page.get("title", ""))

    # 填充副标题
    _fill_subtitle(slide, page.get("subtitle", ""))

    # 填充要点
    bullets = page.get("bullets")
    if bullets:
        _fill_body(slide, bullets)

    # 添加图表
    chart = page.get("chart")
    if chart:
        _add_chart(slide, chart)

    # 填充备注
    _fill_notes(slide, page.get("notes", ""))


def render_pptx(
    data: dict,
    out_path: str,
    template_path: str | None = None,
) -> dict:
    """将结构化大纲渲染为 PPTX 文件。

    Args:
        data: LLM 生成的大纲 JSON
        out_path: 输出文件路径
        template_path: PPT 模板路径（可选）

    Returns:
        包含输出路径和渲染信息的字典
    """
    path = Path(out_path)
    path.parent.mkdir(parents=True, exist_ok=True)

    # 加载模板或创建空白演示文稿
    if template_path and Path(template_path).exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()

    pages = data.get("pages", [])
    for page in pages:
        render_slide(prs, page)

    prs.save(str(path))

    return {
        "out": str(path),
        "page_count": len(pages),
        "template_used": bool(template_path and Path(template_path).exists()),
    }
