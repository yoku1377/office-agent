"""母版解析：扫描 PPT 模板，提取各版式占位符清单（含位置尺寸），作为大纲生成的"版式菜单"。"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Emu


def _emu_to_inches(emu) -> float:
    """将 EMU 单位转换为英寸（保留1位小数）。"""
    if emu is None:
        return 0.0
    return round(Emu(emu).inches, 1)


@dataclass
class PlaceholderInfo:
    idx: int
    name: str
    type_name: str
    type_id: int
    left: float = 0.0    # 英寸
    top: float = 0.0     # 英寸
    width: float = 0.0   # 英寸
    height: float = 0.0  # 英寸


@dataclass
class LayoutInfo:
    name: str
    placeholders: list[PlaceholderInfo] = field(default_factory=list)


# python-pptx 占位符类型映射
_PLACEHOLDER_TYPE_NAMES = {
    PP_PLACEHOLDER.TITLE: "title",
    PP_PLACEHOLDER.BODY: "body",
    PP_PLACEHOLDER.CENTER_TITLE: "center_title",
    PP_PLACEHOLDER.SUBTITLE: "subtitle",
    PP_PLACEHOLDER.OBJECT: "object",
    PP_PLACEHOLDER.CHART: "chart",
    PP_PLACEHOLDER.TABLE: "table",
    PP_PLACEHOLDER.PICTURE: "picture",
    PP_PLACEHOLDER.ORG_CHART: "org_chart",
    PP_PLACEHOLDER.MEDIA_CLIP: "media_clip",
    PP_PLACEHOLDER.SLIDE_NUMBER: "slide_number",
    PP_PLACEHOLDER.HEADER: "header",
    PP_PLACEHOLDER.FOOTER: "footer",
    PP_PLACEHOLDER.DATE: "date",
}


def _placeholder_type_name(ph_type: int) -> str:
    return _PLACEHOLDER_TYPE_NAMES.get(ph_type, f"unknown({ph_type})")


def parse_template(template_path: str | Path) -> list[LayoutInfo]:
    """解析 PPT 模板，返回所有版式及其占位符信息（含位置尺寸）。"""
    path = Path(template_path)
    if not path.exists():
        return []

    prs = Presentation(str(path))
    layouts: list[LayoutInfo] = []

    for layout in prs.slide_layouts:
        info = LayoutInfo(name=layout.name)
        for ph in layout.placeholders:
            info.placeholders.append(
                PlaceholderInfo(
                    idx=ph.placeholder_format.idx,
                    name=ph.name,
                    type_name=_placeholder_type_name(ph.placeholder_format.type),
                    type_id=ph.placeholder_format.type,
                    left=_emu_to_inches(ph.left),
                    top=_emu_to_inches(ph.top),
                    width=_emu_to_inches(ph.width),
                    height=_emu_to_inches(ph.height),
                )
            )
        layouts.append(info)

    return layouts


def build_layout_menu(layouts: list[LayoutInfo]) -> str:
    """将版式信息格式化为 LLM 可读的版式菜单文本（含位置尺寸）。"""
    if not layouts:
        return _default_layout_menu()

    lines = []
    for i, layout in enumerate(layouts, 1):
        ph_parts = []
        for ph in layout.placeholders:
            size_info = f"{ph.width}x{ph.height}in@({ph.left},{ph.top})"
            ph_parts.append(f"{ph.name}(idx={ph.idx}, type={ph.type_name}, {size_info})")
        ph_desc = ", ".join(ph_parts)
        lines.append(f"  {i}. \"{layout.name}\" — 占位符: [{ph_desc}]" if ph_desc
                     else f"  {i}. \"{layout.name}\" — 无占位符")
    return "\n".join(lines)


def _default_layout_menu() -> str:
    """无模板时的默认版式菜单（含典型尺寸）。"""
    return """  1. "Title Slide" — 占位符: [Title 1(idx=0, type=title, 8.0x1.5in@(2.7,1.8)), Subtitle 2(idx=1, type=subtitle, 8.0x1.5in@(2.7,3.8))]
  2. "Title and Content" — 占位符: [Title 1(idx=0, type=title, 11.3x1.0in@(0.5,0.3)), Content Placeholder 2(idx=1, type=body, 11.3x5.5in@(0.5,1.5))]
  3. "Section Header" — 占位符: [Title 1(idx=0, type=title, 11.3x1.0in@(0.5,0.3)), Text Placeholder 2(idx=1, type=body, 11.3x5.5in@(0.5,1.5))]
  4. "Two Content" — 占位符: [Title 1(idx=0, type=title, 11.3x1.0in@(0.5,0.3)), Content Placeholder 2(idx=1, type=body, 5.3x5.5in@(0.5,1.5)), Content Placeholder 3(idx=2, type=body, 5.3x5.5in@(6.8,1.5))]
  5. "Comparison" — 占位符: [Title 1(idx=0, type=title, 11.3x1.0in@(0.5,0.3)), Content Placeholder 2(idx=1, type=body, 5.3x5.5in@(0.5,1.5)), Content Placeholder 3(idx=2, type=body, 5.3x5.5in@(6.8,1.5))]
  6. "Title Only" — 占位符: [Title 1(idx=0, type=title, 11.3x1.0in@(0.5,0.3))]
  7. "Blank" — 无占位符
  8. "Content with Caption" — 占位符: [Title 1(idx=0, type=title, 11.3x1.0in@(0.5,0.3)), Content Placeholder 2(idx=1, type=body, 11.3x4.0in@(0.5,1.5)), Caption Placeholder 3(idx=2, type=body, 11.3x1.0in@(0.5,5.8))]
  9. "Picture with Caption" — 占位符: [Title 1(idx=0, type=title, 11.3x1.0in@(0.5,0.3)), Picture Placeholder 2(idx=1, type=picture, 8.0x4.0in@(2.0,1.5)), Caption Placeholder 3(idx=2, type=body, 11.3x1.0in@(0.5,5.8))]"""


def get_layout_by_name(prs: Presentation, name: str):
    """根据名称查找版式，找不到则回退到第一个版式。"""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    # 回退：尝试模糊匹配
    for layout in prs.slide_layouts:
        if name.lower() in layout.name.lower():
            return layout
    # 最终回退到第一个版式
    return prs.slide_layouts[0] if prs.slide_layouts else None
